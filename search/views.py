import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

from rest_framework import status
from rest_framework.response import Response
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated
from rest_framework.generics import ListAPIView
from rest_framework.parsers import MultiPartParser

from django.conf import settings
from django.shortcuts import render, redirect
from django.contrib.auth import authenticate
from django.contrib.auth.models import User
from django.contrib.auth.tokens import default_token_generator
from django.contrib.auth.password_validation import validate_password
from django.utils.encoding import force_bytes, force_str
from django.utils.http import urlsafe_base64_encode, urlsafe_base64_decode
from django.core.exceptions import ValidationError as DjangoValidationError

from rest_framework_simplejwt.tokens import RefreshToken
from .serializers import RegisterSerializer
from .ML.intent_model import classify_intent, map_intent_to_model

from .models import UserProfile, SearchQuery
from .serializers import (
    SearchQuerySerializer,
    UserProfileSerializer,
    ChangePasswordSerializer,
)
from .services.perplexity import call_perplexity_model, call_groq_model
from .services.chat_context import build_chat_context
from .utils import (
    image_to_data_uri,
    send_verification_email,
    send_password_reset_email,
    get_best_model,
)


FRONTEND_BASE_URL = settings.FRONTEND_BASE_URL
BACKEND_BASE_URL = settings.BACKEND_BASE_URL


def test_ui_view(request):
    return render(request, "index.html")


class RegisterView(APIView):
    def post(self, request):
        serializer = RegisterSerializer(data=request.data)

        if serializer.is_valid():
            user = serializer.save()

            uid = urlsafe_base64_encode(force_bytes(user.pk))
            token = default_token_generator.make_token(user)

            backend_base_url = f"{BACKEND_BASE_URL}/verify-email"
            verification_url = f"{backend_base_url}/{uid}/{token}/"

            send_verification_email(user, verification_url)

            return Response(
                {"message": "User created. Check your email to verify."}, status=201
            )

        return Response(serializer.errors, status=400)


class ChangePasswordView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        serializer = ChangePasswordSerializer(data=request.data)
        user = request.user

        if serializer.is_valid():
            if not user.check_password(serializer.validated_data["current_password"]):
                return Response(
                    {"current_password": ["Incorrect current password."]},
                    status=status.HTTP_400_BAD_REQUEST,
                )

            user.set_password(serializer.validated_data["new_password"])
            user.save()
            return Response(
                {"detail": "Password changed successfully."}, status=status.HTTP_200_OK
            )

        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


class UserProfileView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        serializer = UserProfileSerializer(request.user.userprofile)
        return Response(serializer.data)


class UserProfileUpdateView(APIView):
    permission_classes = [IsAuthenticated]

    def patch(self, request):
        profile = UserProfile.objects.get(user=request.user)
        serializer = UserProfileSerializer(profile, data=request.data, partial=True)
        if serializer.is_valid():
            serializer.save()
            return Response(serializer.data)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


class VerifyEmailView(APIView):
    def get(self, request, uidb64, token):
        try:
            uid = urlsafe_base64_decode(uidb64).decode()
            user = User.objects.get(pk=uid)
        except (User.DoesNotExist, ValueError, TypeError, OverflowError):
            return redirect(f"{FRONTEND_BASE_URL}?status=invalid")

        if default_token_generator.check_token(user, token):
            profile = UserProfile.objects.get(user=user)
            profile.is_verified = True
            profile.save()
            return redirect(f"{FRONTEND_BASE_URL}?status=success")

        return redirect(f"{FRONTEND_BASE_URL}?status=expired")


class RequestPasswordResetView(APIView):
    def post(self, request):
        email = request.data.get("email")
        if not email:
            return Response({"error": "Email is required"}, status=400)

        try:
            user = User.objects.get(email=email)
            uid = urlsafe_base64_encode(force_bytes(user.pk))
            token = default_token_generator.make_token(user)

            reset_link = (
                f"{FRONTEND_BASE_URL}?event=reset-password&uidb64={uid}&token={token}"
            )

            send_password_reset_email(user, reset_link)

            return Response(
                {"message": "Password reset link sent to your email"}, status=200
            )

        except User.DoesNotExist:
            return Response(
                {"error": "User with this email does not exist"}, status=404
            )


class PasswordResetConfirmView(APIView):
    def post(self, request, uidb64, token):
        try:
            uid = force_str(urlsafe_base64_decode(uidb64))
            user = User.objects.get(pk=uid)
        except (User.DoesNotExist, ValueError, TypeError, OverflowError):
            return Response({"error": "Invalid or expired link"}, status=400)

        if not default_token_generator.check_token(user, token):
            return Response({"error": "Invalid or expired token"}, status=400)

        new_password = request.data.get("new_password")
        if not new_password:
            return Response({"error": "New password is required"}, status=400)

        try:
            validate_password(new_password)
        except DjangoValidationError as e:
            return Response({"error": e.messages}, status=400)

        user.set_password(new_password)
        user.save()

        return Response({"reset-password-status": "success"})


class LoginView(APIView):
    def post(self, request):
        username = request.data.get("username")
        password = request.data.get("password")
        user = authenticate(username=username, password=password)

        if user is not None:
            profile = UserProfile.objects.filter(user=user).first()
            if not profile or not profile.is_verified:
                return Response({"error": "Email not verified"}, status=403)

            refresh = RefreshToken.for_user(user)
            return Response(
                {
                    "access": str(refresh.access_token),
                    "refresh": str(refresh),
                }
            )

        return Response({"error": "Invalid credentials"}, status=401)


class LogoutView(APIView):
    def post(self, request):
        try:
            refresh_token = request.data["refresh"]
            token = RefreshToken(refresh_token)
            token.blacklist()
            return Response(status=205)
        except Exception as e:
            return Response(status=400)


class SearchView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        prompt = request.data.get("prompt")
        image_url = request.data.get("image_url")
        # model = request.data.get("model", "sonar")
        # return_images = request.data.get("return_images", False)
        return_images = False
        search_mode = request.data.get("search_mode", "web")
        deep_research = request.data.get("deep_research", False)
        pro = request.data.get("pro", True)
        labs = request.data.get("labs", False)

        # model = get_best_model(model)

        if prompt:
            intent, confidence = classify_intent(prompt)
            model = map_intent_to_model(intent)

            # BUILD PROMPT WITH PREVIOUS MESSSAGES (CHAT HISTORY) FOR PROVIDING CONTEXT TO THE MODEL
            final_prompt = build_chat_context(request.user, prompt)

        else:
            model = "sonar-pro"

        if not prompt and not image_url:
            return Response({"error": "Prompt or image is required."}, status=400)

        # Check if user is subscribed
        profile = UserProfile.objects.filter(user=request.user).first()
        if not profile or not profile.is_subscribed:
            # Apply search limit for non-subscribed users
            search_count = SearchQuery.objects.filter(user=request.user).count()
            if search_count >= int(os.getenv("FREE_SEARCH_LIMIT", 10)):
                return Response(
                    {"error": "Free search limit reached. Please subscribe."},
                    status=402,
                )

        result = call_perplexity_model(
            prompt=final_prompt,
            image_url=image_url,
            model=model,
            return_images=return_images,
            search_mode=search_mode,
            deep_research=deep_research,
        )

        # result = call_groq_model(
        #     prompt=prompt,
        #     image_url=image_url,
        #     model=model,
        #     return_images=return_images,
        #     search_mode=search_mode,
        #     deep_research=deep_research,
        # )

        SearchQuery.objects.create(
            user=request.user, prompt=prompt or "[Image]", response=result
        )

        return Response(result)


# library functionality


class LibraryView(ListAPIView):
    serializer_class = SearchQuerySerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        return SearchQuery.objects.filter(user=self.request.user).order_by(
            "-created_at"
        )[:10]


class UploadImageView(APIView):
    parser_classes = [MultiPartParser]

    def post(self, request):
        image = request.FILES.get("image")
        if not image:
            return Response({"error": "No image uploaded"}, status=400)

        try:
            data_uri = image_to_data_uri(image)
            return Response({"image_url": data_uri}, status=200)
        except Exception as e:
            return Response({"error": str(e)}, status=500)


class UploadDocExtractView(APIView):
    parser_classes = [MultiPartParser]

    def post(self, request):
        uploaded_file = request.FILES.get("file")
        if not uploaded_file:
            return Response({"error": "No file uploaded"}, status=400)

        try:
            text = self.extract_text(uploaded_file)
            return Response({"text_content": text}, status=200)
        except Exception as e:
            return Response({"error": str(e)}, status=500)

    # def extract_text(self, file):
    #     if file.name.endswith(".txt"):
    #         return file.read().decode("utf-8")

    #     elif file.name.endswith(".pdf"):
    #         reader = PdfReader(file)
    #         text = "\n".join(
    #             [page.extract_text() for page in reader.pages if page.extract_text()]
    #         )
    #         return text.strip()

    #     elif file.name.endswith(".docx"):
    #         doc = Document(file)
    #         return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

    #     else:
    #         raise ValueError(
    #             "Unsupported file format. Please upload a .txt, .pdf, or .docx file."
    #         )


    def extract_text(self, file):
        if file.name.endswith(".txt"):
            return file.read().decode("utf-8")

        elif file.name.endswith(".pdf"):
            reader = PdfReader(file)
            text = "\n".join(
                [page.extract_text() for page in reader.pages if page.extract_text()]
            )
            return text.strip()

        elif file.name.endswith(".docx"):
            doc = Document(file)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif file.name.endswith(".csv"):
            df = pd.read_csv(file)
            return df.to_markdown(index=False)

        elif file.name.endswith(".xlsx"):
            # Read all sheets
            excel_file = pd.read_excel(file, sheet_name=None)
            extracted = []
            for sheet_name, df in excel_file.items():
                extracted.append(f"🧾 **Sheet: {sheet_name}**\n")
                extracted.append(df.to_markdown(index=False))
                extracted.append("\n")
            return "\n".join(extracted)

        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            extracted = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                notes = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                extracted.append(f"🖼️ Slide {i}")
                if slide_text:
                    extracted.append("Content: " + " | ".join(slide_text))
                if notes:
                    extracted.append("Speaker Notes: " + notes)
                extracted.append("\n")
            return "\n".join(extracted).strip()

        else:
            raise ValueError(
                "Unsupported file format. Please upload a .txt, .pdf, .docx, .csv, .xlsx, or .pptx file."
            )
